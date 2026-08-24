from __future__ import annotations

import html
import json
from pathlib import Path
from typing import Any

from neoag.agent_skills.appm_review import main as appm_review_main
from neoag.agent_skills.ccf_review import main as ccf_review_main
from neoag.controlled_execution.io_utils import load_limited_yaml, markdown_table, read_tsv, write_json, write_tsv
from neoag.reports_dual import load_report_bundle, make_patient_report
from neoag.skill_taxonomy.review_skills import (
    run_concept_explainer,
    run_experiment_design,
    run_ranking_compare,
    run_technical_report,
)

from .contracts import MacroResult, MacroStep
from .errors import FailureCode
from .execution_adapters import discover_result_artifacts
from .review_integrity import audit_review_inputs
from .state import RunLayout, audit, new_run_id, safe_identifier, update_case_state


def _get(row: dict[str, str], *keys: str, default: str = "") -> str:
    for key in keys:
        value = row.get(key)
        if value not in {None, ""}:
            return str(value)
    return default


def _truthy(value: Any) -> bool:
    return str(value or "").strip().lower() in {"1", "true", "yes", "y"}


def _integer(value: Any, default: int = 10**9) -> int:
    try:
        return int(float(str(value)))
    except (TypeError, ValueError):
        return default


def _event_kind(event_type: str, consequence: str = "") -> str:
    # The explicit biological event type is authoritative. A DNA SNV/InDel
    # can have a splice-related VEP consequence without becoming an RNA
    # splice-junction event.
    explicit = str(event_type or "").strip().lower()
    explicit_kinds = {
        "snv": "MISSENSE",
        "missense": "MISSENSE",
        "substitution": "MISSENSE",
        "indel": "FRAMESHIFT",
        "insertion": "FRAMESHIFT",
        "deletion": "FRAMESHIFT",
        "frameshift": "FRAMESHIFT",
        "fusion": "FUSION",
        "splice": "SPLICE",
        "splice_junction": "SPLICE",
        "junction": "SPLICE",
        "dna_sv": "DNA_SV",
        "structural_variant": "DNA_SV",
    }
    if explicit in explicit_kinds:
        return explicit_kinds[explicit]

    text = f"{event_type} {consequence}".lower()
    if "fusion" in text:
        return "FUSION"
    if any(token in text for token in ("splice", "junction", "exon")):
        return "SPLICE"
    if any(token in text for token in ("frameshift", "frame_shift", "novel_tail")):
        return "FRAMESHIFT"
    if any(token in text for token in ("structural", "dna_sv", " bnd", " sv")):
        return "DNA_SV"
    if any(token in text for token in ("snv", "missense")):
        return "MISSENSE"
    return "OTHER"


def _recommended_validation(kind: str) -> str:
    return {
        "MISSENSE": "MT/WT paired short peptide + ELISpot; tetramer for confirmed responses",
        "FRAMESHIFT": "novel-tail long peptide and/or frameshift minigene",
        "FUSION": "RT-PCR/Sanger first; fusion-junction long peptide and/or minigene",
        "SPLICE": "targeted RNA first; abnormal-junction long peptide and/or splice minigene",
        "DNA_SV": "DNA breakpoint + RNA transcript confirmation before SV minigene",
    }.get(kind, "manual evidence review before assay selection")


def _priority(grade: str, row: dict[str, str]) -> tuple[str, str, str]:
    grade = grade.upper() or "R4"
    reasons: list[str] = []
    hard = _truthy(row.get("hard_failure")) or bool(_get(row, "hard_failure_codes"))
    manual = _truthy(row.get("manual_review_required"))
    kind = _event_kind(_get(row, "event_type", "evidence_track"), _get(row, "consequence"))
    rna = _get(row, "rna_support_state", "rna_support_status", default="RNA_UNASSESSED").upper()
    safety = _get(row, "safety_state", "safety_status", default="SAFETY_PARTIAL").upper()
    presentation = _get(row, "presentation_consensus_state", "presentation_state", default="PRESENTATION_UNASSESSED").upper()
    phase = _get(row, "phase_status", "haplotype_status", "phase_confidence").upper()
    fusion = _get(row, "fusion_consensus_status", "event_authenticity_state", "tools_detected").upper()
    source_chain_tier = _get(row, "source_chain_confidence_tier").upper()

    if source_chain_tier == "C4":
        return "EXCLUDE_HARD_FAIL", "NOT_FOR_EXPERIMENT", "REVIEW_SOURCE_CHAIN_C4"
    if hard:
        return "EXCLUDE_HARD_FAIL", "NOT_FOR_EXPERIMENT", "REVIEW_HARD_FAIL"
    if grade == "R4":
        return ("MANUAL_REVIEW", "MANUAL_REVIEW_ONLY", "REVIEW_DRIVER_OR_MECHANISM") if manual else ("HOLD", "HOLD", "REVIEW_R4_HOLD")
    if source_chain_tier == "C3":
        return "COMPLETE_EVIDENCE", "SOURCE_CHAIN_COMPLETION_FIRST", "REVIEW_SOURCE_CHAIN_C3_INCOMPLETE"
    if "PHAS" in phase and any(token in phase for token in ("REQUIRED", "UNRESOLVED", "LOW")):
        return "COMPLETE_EVIDENCE", "PHASING_FIRST", "REVIEW_PHASING_REQUIRED"
    if kind == "FUSION" and ("SINGLE" in fusion or "ONE_CALLER" in fusion):
        return "COMPLETE_EVIDENCE", "FUSION_CONFIRMATION_FIRST", "REVIEW_FUSION_SINGLE_CALLER"
    if any(token in rna for token in ("UNASSESSED", "MISSING", "NOT_DETECTED", "NEGATIVE", "GENE_EXPRESSION_ONLY")):
        return "COMPLETE_EVIDENCE", "TARGETED_RNA_FIRST", "REVIEW_RNA_CONFIRMATION_REQUIRED"
    if any(token in safety for token in ("PARTIAL", "REVIEW", "UNASSESSED", "CAUTION", "HIGH_RISK")):
        return "COMPLETE_EVIDENCE", "SAFETY_REVIEW_FIRST", "REVIEW_SAFETY_INCOMPLETE"
    if any(token in presentation for token in ("SINGLE_TOOL", "DISCORDANT", "UNASSESSED")):
        return "COMPLETE_EVIDENCE", "PRESENTATION_REVIEW_FIRST", "REVIEW_PRESENTATION_INCOMPLETE"
    if grade == "R1":
        reasons.append("REVIEW_R1_COMPLETE_CORE_EVIDENCE")
        return "ADVANCE_EXPERIMENT", "EXPERIMENT_PRIORITY_HIGH", ";".join(reasons)
    if grade == "R2":
        reasons.append("REVIEW_R2_ONE_OR_MORE_CAUTIONS")
        return "ADVANCE_EXPERIMENT", "EXPERIMENT_PRIORITY_MEDIUM_HIGH", ";".join(reasons)
    return "COMPLETE_EVIDENCE", "EVIDENCE_COMPLETION_FIRST", "REVIEW_R3_COMPLETE_EVIDENCE"


def _representative(event: dict[str, str], peptide_by_id: dict[str, dict[str, str]], all_tool_by_id: dict[str, dict[str, str]], index: int) -> dict[str, str]:
    peptide_id = _get(event, f"representative_{index}_peptide_id", "best_peptide_id" if index == 1 else "")
    merged = dict(event)
    if peptide_id:
        merged.update(peptide_by_id.get(peptide_id, {}))
        merged.update(all_tool_by_id.get(peptide_id, {}))
    merged["peptide_id"] = peptide_id
    merged["peptide"] = _get(event, f"representative_{index}_peptide", default=_get(merged, "peptide"))
    merged["hla_allele"] = _get(event, f"representative_{index}_hla_allele", default=_get(merged, "hla_allele"))
    return merged


def build_review_rows(events: list[dict[str, str]], peptides: list[dict[str, str]], all_tool: list[dict[str, str]]) -> list[dict[str, str]]:
    peptide_by_id = {_get(row, "peptide_id"): row for row in peptides if _get(row, "peptide_id")}
    all_tool_by_id = {_get(row, "peptide_id"): row for row in all_tool if _get(row, "peptide_id")}
    output: list[dict[str, str]] = []
    for event in events:
        rep1 = _representative(event, peptide_by_id, all_tool_by_id, 1)
        rep2 = _representative(event, peptide_by_id, all_tool_by_id, 2)
        grade = _get(event, "best_evidence_grade", "evidence_grade", default=_get(rep1, "evidence_grade", default="R4"))
        merged = {**event, **rep1}
        status, priority, reason = _priority(grade, merged)
        event_type = _get(event, "event_type", "biological_event_track", "evidence_track", default=_get(rep1, "event_type"))
        kind = _event_kind(event_type, _get(rep1, "consequence"))
        output.append({
            "pipeline_event_rank": _get(event, "event_evidence_rank"),
            "pipeline_r_grade": grade,
            "event_group_id": _get(event, "event_group_id", "event_id"),
            "phase_group_id": _get(event, "phase_group_id", default=_get(rep1, "phase_group_id")),
            "redundancy_group": _get(event, "redundancy_group", default=_get(rep1, "redundancy_group")),
            "event_id": _get(event, "event_id", default=_get(rep1, "event_id")),
            "gene": _get(event, "gene", default=_get(rep1, "gene")),
            "event_type": event_type,
            "event_kind": kind,
            "source_chain_track": _get(event, "source_chain_track", default=_get(rep1, "source_chain_track")),
            "source_chain_confidence_tier": _get(event, "source_chain_confidence_tier", default=_get(rep1, "source_chain_confidence_tier")),
            "source_chain_confidence_label": _get(event, "source_chain_confidence_label", default=_get(rep1, "source_chain_confidence_label")),
            "source_chain_orthogonal_status": _get(event, "source_chain_orthogonal_status", default=_get(rep1, "source_chain_orthogonal_status")),
            "source_chain_orthogonal_sources": _get(event, "source_chain_orthogonal_sources", default=_get(rep1, "source_chain_orthogonal_sources")),
            "source_chain_reason_codes": _get(event, "source_chain_reason_codes", default=_get(rep1, "source_chain_reason_codes")),
            "source_chain_missing_requirements": _get(event, "source_chain_missing_requirements", default=_get(rep1, "source_chain_missing_requirements")),
            "source_chain_low_power_requirements": _get(event, "source_chain_low_power_requirements", default=_get(rep1, "source_chain_low_power_requirements")),
            "source_chain_negative_requirements": _get(event, "source_chain_negative_requirements", default=_get(rep1, "source_chain_negative_requirements")),
            "source_chain_hard_failure_codes": _get(event, "source_chain_hard_failure_codes", default=_get(rep1, "source_chain_hard_failure_codes")),
            "pareto_front": _get(event, "best_pareto_front", default=_get(rep1, "pareto_front")),
            "representative_1_peptide_id": _get(rep1, "peptide_id"),
            "representative_1_peptide": _get(rep1, "peptide"),
            "representative_1_hla": _get(rep1, "hla_allele"),
            "representative_2_peptide_id": _get(rep2, "peptide_id"),
            "representative_2_peptide": _get(rep2, "peptide"),
            "representative_2_hla": _get(rep2, "hla_allele"),
            "event_authenticity_state": _get(rep1, "event_authenticity_state", "cross_platform_status"),
            "rna_support_state": _get(rep1, "rna_support_state", "rna_support_status", default="RNA_UNASSESSED"),
            "presentation_consensus_state": _get(rep1, "presentation_consensus_state", "presentation_state"),
            "mutant_specificity_state": _get(rep1, "mutant_specificity_state", "mutant_specificity_status"),
            "clonality_state": _get(rep1, "clonality_state", "clonality_status"),
            "ccf_estimate": _get(rep1, "ccf_estimate", "ccf_best"),
            "ccf_confidence": _get(rep1, "ccf_confidence"),
            "hla_appm_state": _get(rep1, "hla_appm_state", "appm_integrity_status", "escape_status"),
            "safety_state": _get(rep1, "safety_state", "safety_status", default="SAFETY_PARTIAL"),
            "evidence_completeness_state": _get(rep1, "evidence_completeness_state", "appm_evidence_completeness"),
            "evidence_conflict_fields": _get(rep1, "evidence_conflict_fields"),
            "hard_failure": _get(rep1, "hard_failure", default="no"),
            "hard_failure_codes": _get(rep1, "hard_failure_codes"),
            "priority_cap": _get(rep1, "priority_cap"),
            "priority_cap_reason_codes": _get(rep1, "priority_cap_reason_codes"),
            "manual_review_required": _get(event, "manual_review_required", default=_get(rep1, "manual_review_required", default="no")),
            "review_status": status,
            "review_reason": reason,
            "experiment_priority": priority,
            "recommended_validation": _recommended_validation(kind),
            "pipeline_recommended_next_steps": _get(event, "recommended_next_steps", default=_get(rep1, "recommended_next_steps")),
            "pipeline_consensus_action": _get(event, "consensus_action", default=_get(rep1, "consensus_action")),
            "pipeline_consensus_trace": _get(event, "event_consensus_trace", default=_get(rep1, "consensus_trace")),
            "member_event_count": _get(event, "member_event_count"),
            "peptide_count": _get(event, "peptide_count"),
        })
    return output


def select_first_batch(rows: list[dict[str, str]], top_n: int) -> list[dict[str, str]]:
    priority_order = {
        "EXPERIMENT_PRIORITY_HIGH": 0,
        "EXPERIMENT_PRIORITY_MEDIUM_HIGH": 1,
        "FUSION_CONFIRMATION_FIRST": 2,
        "TARGETED_RNA_FIRST": 3,
        "PHASING_FIRST": 4,
        "SOURCE_CHAIN_COMPLETION_FIRST": 5,
        "EVIDENCE_COMPLETION_FIRST": 6,
    }
    clonality_order = {"CLONAL_LIKE": 0, "CLONAL": 0, "SUBCLONAL_LIKE": 1, "SUBCLONAL": 1, "UNRESOLVED": 2, "": 3}
    eligible = [row for row in rows if row["experiment_priority"] in priority_order]
    eligible.sort(key=lambda row: (priority_order[row["experiment_priority"]], clonality_order.get(str(row.get("clonality_state") or "").upper(), 3), _integer(row.get("pipeline_event_rank")), row.get("event_group_id", "")))
    quotas = {"MISSENSE": 5, "FRAMESHIFT": 2, "FUSION": 2, "SPLICE": 2, "DNA_SV": 1, "OTHER": 1}
    selected: list[dict[str, str]] = []
    seen_groups: set[str] = set()
    seen_phase: set[str] = set()
    seen_redundancy: set[str] = set()
    hla_counts: dict[str, int] = {}

    def can_add(row: dict[str, str]) -> bool:
        group = row.get("event_group_id") or row.get("event_id")
        phase = row.get("phase_group_id")
        redundancy = row.get("redundancy_group")
        hla = row.get("representative_1_hla")
        return bool(group and group not in seen_groups and (not phase or phase not in seen_phase) and (not redundancy or redundancy not in seen_redundancy) and (not hla or hla_counts.get(hla, 0) < max(2, top_n // 3)))

    def add(row: dict[str, str], method: str) -> None:
        copy = dict(row)
        copy["selection_method"] = method
        selected.append(copy)
        seen_groups.add(row.get("event_group_id") or row.get("event_id", ""))
        if row.get("phase_group_id"):
            seen_phase.add(row["phase_group_id"])
        if row.get("redundancy_group"):
            seen_redundancy.add(row["redundancy_group"])
        hla = row.get("representative_1_hla")
        if hla:
            hla_counts[hla] = hla_counts.get(hla, 0) + 1

    for kind, quota in quotas.items():
        for row in [candidate for candidate in eligible if candidate.get("event_kind") == kind]:
            if sum(item.get("event_kind") == kind for item in selected) >= quota or len(selected) >= top_n:
                break
            if can_add(row):
                add(row, "deterministic_grade_event_phase_redundancy_hla_type_diversity")
    for row in eligible:
        if len(selected) >= top_n:
            break
        if can_add(row):
            add(row, "deterministic_fill_after_type_quotas")
    manual = [row for row in rows if row["experiment_priority"] == "MANUAL_REVIEW_ONLY"]
    if manual and len(selected) < top_n and can_add(manual[0]):
        add(manual[0], "single_manual_review_lane_not_auto_promoted")
    for index, row in enumerate(selected, 1):
        row["first_batch_rank"] = str(index)
    return selected


def _read_context(*paths: str | Path | None) -> dict[str, Any]:
    merged: dict[str, Any] = {}
    for value in paths:
        if not value or not Path(value).is_file():
            continue
        try:
            data = load_limited_yaml(value)
        except Exception:
            continue
        if isinstance(data, dict):
            merged.update(data)
    return merged


def _context_summary(context: dict[str, Any]) -> str:
    keys = ("disease", "diagnosis", "sample_type", "analysis_goal", "therapy_context", "notes")
    rows = [f"- {key}: {context[key]}" for key in keys if str(context.get(key, "")).strip()]
    return "\n".join(rows) if rows else "No structured clinical or disease context was supplied."


def _read_preview(path: str | Path | None, limit: int = 12000) -> str:
    if not path or not Path(path).is_file():
        return "UNASSESSED: supporting report was not available."
    return Path(path).read_text(encoding="utf-8", errors="replace")[:limit]


def _write_docx(path: Path, title: str, sections: list[tuple[str, str]], rows: list[dict[str, str]]) -> bool:
    try:
        from docx import Document
    except Exception:
        return False
    doc = Document()
    doc.add_heading(title, 0)
    doc.add_paragraph("研究性计算筛选结果；不构成临床诊断、治疗建议或疗效承诺。")
    for heading, body in sections:
        doc.add_heading(heading, level=1)
        doc.add_paragraph(body)
    cols = ["first_batch_rank", "gene", "event_kind", "pipeline_r_grade", "experiment_priority", "recommended_validation"]
    table = doc.add_table(rows=1, cols=len(cols)); table.style = "Table Grid"
    for i, col in enumerate(cols): table.rows[0].cells[i].text = col
    for row in rows:
        cells = table.add_row().cells
        for i, col in enumerate(cols): cells[i].text = str(row.get(col, ""))
    doc.save(path)
    return True


def _count_data_rows(path: Path) -> int:
    if not path.is_file():
        return 0
    try:
        with path.open("r", encoding="utf-8", errors="replace") as handle:
            return max(0, sum(1 for line in handle if line.strip()) - 1)
    except OSError:
        return 0


def _first_existing_path(paths: list[Path]) -> Path | None:
    for path in paths:
        if path.is_file():
            return path
    return None


def _fusion_root(result_dir: Path) -> Path | None:
    for candidate in (
        result_dir / "pipeline" / "production" / "branches" / "fusion",
        result_dir / "branches" / "fusion",
        result_dir / "fusion",
        result_dir.parent / "branches" / "fusion",
        result_dir.parent.parent / "branches" / "fusion",
    ):
        if candidate.exists():
            return candidate
    return None


def _fusion_report_data(result_dir: Path) -> tuple[list[dict[str, str]], list[dict[str, str]], str]:
    root = _fusion_root(result_dir)
    if root is None:
        return [], [], "No fusion branch directory was found in the reviewed result directory."
    rows: list[dict[str, str]] = []

    def add(tool: str, path: Path | None, note: str = "") -> None:
        rows.append({
            "tool_or_layer": tool,
            "records": str(_count_data_rows(path) if path else 0),
            "path": str(path) if path and path.exists() else "",
            "note": note,
        })

    easyfuse_files = sorted(root.glob("easyfuse/*/fusions.pass.csv")) + sorted(root.glob("easyfuse/fusions.pass.csv"))
    easyfuse_count = sum(_count_data_rows(path) for path in easyfuse_files)
    rows.append({
        "tool_or_layer": "EasyFuse PASS calls",
        "records": str(easyfuse_count),
        "path": ";".join(str(path) for path in easyfuse_files),
        "note": "caller-level PASS fusion calls before Open-Neo peptide-HLA generation",
    })
    add("STAR-Fusion calls", _first_existing_path([root / "star-fusion" / "star-fusion.fusion_predictions.tsv", root / "star_fusion" / "star-fusion.fusion_predictions.tsv"]), "optional caller")
    add("Arriba calls", _first_existing_path(sorted(root.glob("arriba/*.fusions.tsv")) + [root / "arriba" / "fusions.tsv"]), "optional caller")
    add("FusionCatcher calls", _first_existing_path([root / "fusioncatcher" / "fusioncatcher.final-list.txt", root / "fusioncatcher" / "final-list_candidate-fusion-genes.txt"]), "optional caller")
    add("Fusion consensus", root / "consensus" / "fusion_consensus.tsv", "cross-caller consensus layer")
    raw_events_path = root / "intermediates" / "parsed" / "raw_events.tsv"
    raw_peptides_path = root / "intermediates" / "parsed" / "raw_peptides.tsv"
    add("Open-Neo fusion raw events", raw_events_path, "normalized fusion events forwarded to the pipeline")
    add("Open-Neo fusion peptides", raw_peptides_path, "fusion peptide-HLA rows entering presentation/ranking")

    top_events: list[dict[str, str]] = []
    if raw_events_path.is_file():
        _, event_rows = read_tsv(raw_events_path)
        for event in event_rows[:20]:
            top_events.append({
                "fusion": _get(event, "gene", "event_name"),
                "frame": _get(event, "rna_frame_status", "frame_status", "consequence"),
                "junction_reads": _get(event, "rna_junction_reads", "provided_rna_junction_reads"),
                "confidence": _get(event, "event_confidence", "confidence"),
                "source": _get(event, "source", "source_tools", "source_file"),
            })
    peptide_count = _count_data_rows(raw_peptides_path)
    event_count = _count_data_rows(raw_events_path)
    if peptide_count == 0 and (easyfuse_count or event_count):
        note = "Fusion detection produced caller/event records, but no fusion peptide-HLA row entered the ranked neoantigen table for this run. Report this as no reportable fusion peptide in the current run, not as fusion not run."
    elif peptide_count:
        note = "Fusion peptide-HLA rows were generated and are eligible for downstream presentation/ranking review."
    else:
        note = "No reportable fusion calls were found in the reviewed result directory."
    return rows, top_events, note


def _fusion_report_markdown(result_dir: Path, out_tsv: Path | None = None) -> str:
    rows, top_events, note = _fusion_report_data(result_dir)
    if out_tsv is not None and rows:
        write_tsv(out_tsv, rows)
    parts = [note]
    if rows:
        parts += ["", markdown_table(rows, max_rows=50)]
    if top_events:
        parts += ["", "Top normalized fusion events:", "", markdown_table(top_events, max_rows=20)]
    return "\n".join(parts)


def _fusion_report_html(result_dir: Path) -> str:
    rows, top_events, note = _fusion_report_data(result_dir)
    if not rows:
        return ""

    def table(records: list[dict[str, str]], keys: list[str]) -> str:
        head = "".join(f"<th>{html.escape(key)}</th>" for key in keys)
        body = []
        for record in records:
            body.append("<tr>" + "".join(f"<td>{html.escape(str(record.get(key, '')))}</td>" for key in keys) + "</tr>")
        return "<table><thead><tr>" + head + "</tr></thead><tbody>" + "".join(body) + "</tbody></table>"

    content = ["<div class='section'><h2>融合基因检测结果</h2>", f"<p>{html.escape(note)}</p>"]
    content.append(table(rows, ["tool_or_layer", "records", "note", "path"]))
    if top_events:
        content.append("<h3>代表性 fusion event</h3>")
        content.append(table(top_events, ["fusion", "frame", "junction_reads", "confidence", "source"]))
    content.append("</div>")
    return "\n".join(content)


def _write_onepage(path: Path, first_batch: list[dict[str, str]], integrity: dict[str, Any]) -> bool:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception:
        return False
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.6))
    title.text_frame.text = "Open-Neo event-level experimental review"
    title.text_frame.paragraphs[0].font.size = Pt(26); title.text_frame.paragraphs[0].font.bold = True
    summary = slide.shapes.add_textbox(Inches(0.65), Inches(1.15), Inches(4.0), Inches(5.7))
    summary.text_frame.text = f"Integrity: {integrity.get('status', 'UNASSESSED')}\nFirst batch events: {len(first_batch)}\n\nResearch boundary\nComputational candidates require experimental validation."
    body = slide.shapes.add_textbox(Inches(4.8), Inches(1.15), Inches(7.9), Inches(5.7))
    lines = [f"{row.get('first_batch_rank')}. {row.get('gene')} | {row.get('event_kind')} | {row.get('pipeline_r_grade')} | {row.get('experiment_priority')}" for row in first_batch[:12]]
    body.text_frame.text = "First-batch event set\n\n" + ("\n".join(lines) if lines else "No event met first-batch criteria.")
    prs.save(path)
    return True


def _write_reports(layout: RunLayout, context: dict[str, Any], review_rows: list[dict[str, str]], first_batch: list[dict[str, str]], artifacts: dict[str, str], integrity: dict[str, Any], support_outputs: dict[str, str], reports: set[str], result_dir: Path) -> dict[str, str]:
    uncertainties = sum(row["review_status"] == "COMPLETE_EVIDENCE" for row in review_rows)
    outputs: dict[str, str] = {}
    technical_sections = [
        ("Integrity and provenance", markdown_table(integrity.get("checks", []), max_rows=100)),
        ("Source artifacts", markdown_table([{"artifact": key, "path": value} for key, value in sorted(artifacts.items())], max_rows=100)),
        ("Fusion detection summary", _fusion_report_markdown(result_dir, layout.review / "fusion_summary.tsv")),
        ("Weighted versus Evidence consensus", _read_preview(support_outputs.get("ranking_compare_report"))),
        ("HLA LOH / APPM review", _read_preview(support_outputs.get("appm_review"))),
        ("CCF / clonality review", _read_preview(support_outputs.get("ccf_review"))),
        ("Conflicts and reason codes", markdown_table([row for row in review_rows if row.get("evidence_conflict_fields") or row.get("hard_failure_codes") or row.get("priority_cap_reason_codes")], max_rows=100)),
        ("Evidence-domain event review", markdown_table(review_rows, max_rows=100)),
        ("Interpretation boundary", "Pipeline R grades and event ranks are preserved. Review status and experiment priority are independent fields. Missing evidence is not a negative result."),
    ]
    if "technical" in reports:
        technical_md = ["# Open-Neo technical review report", ""]
        for heading, body in technical_sections: technical_md += [f"## {heading}", "", body, ""]
        technical_path = layout.reports / "technical_report.md"; technical_path.write_text("\n".join(technical_md) + "\n", encoding="utf-8")
        technical_html = layout.reports / "technical_report.html"; technical_html.write_text("<html><body><pre>" + html.escape("\n".join(technical_md)) + "</pre></body></html>", encoding="utf-8")
        outputs.update({"technical_report_md": str(technical_path), "technical_report_html": str(technical_html)})
        technical_docx = layout.reports / "technical_report.docx"
        if _write_docx(technical_docx, "Open-Neo Technical Review", technical_sections, first_batch): outputs["technical_report_docx"] = str(technical_docx)
    if "onepage" in reports:
        onepage = layout.reports / "onepage_summary.pptx"
        if _write_onepage(onepage, first_batch, integrity): outputs["onepage_summary_pptx"] = str(onepage)
    return outputs


def _read_json_mapping(path: str | Path | None) -> dict[str, Any]:
    if not path or not Path(path).is_file():
        return {}
    try:
        payload = json.loads(Path(path).read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {}
    return payload if isinstance(payload, dict) else {}


def _summary_mapping(path: str | Path | None) -> dict[str, Any]:
    if not path or not Path(path).is_file():
        return {}
    _, rows = read_tsv(path)
    if not rows:
        return {}
    if {"field", "value"}.issubset(rows[0]):
        return {str(row.get("field") or ""): row.get("value", "") for row in rows if row.get("field")}
    if {"metric", "value"}.issubset(rows[0]):
        return {str(row.get("metric") or ""): row.get("value", "") for row in rows if row.get("metric")}
    return dict(rows[0])


def _formal_patient_report(
    layout: RunLayout,
    result_dir: Path,
    context: dict[str, Any],
    artifacts: dict[str, str],
    events: list[dict[str, str]],
    peptides: list[dict[str, str]],
    all_tool: list[dict[str, str]],
) -> dict[str, str]:
    all_tool_by_id = {_get(row, "peptide_id"): row for row in all_tool if _get(row, "peptide_id")}
    report_peptides: list[dict[str, str]] = []
    for peptide in peptides:
        merged = dict(all_tool_by_id.get(_get(peptide, "peptide_id"), {}))
        merged.update(peptide)
        report_peptides.append(merged)

    provenance = _read_json_mapping(artifacts.get("run_manifest"))
    provenance.update(_read_json_mapping(artifacts.get("provenance")))
    provenance["report_role"] = "final_review"
    for key, value in context.items():
        provenance.setdefault(key, value)
    provenance.setdefault("parallel_rankings", {})
    if isinstance(provenance["parallel_rankings"], dict):
        provenance["parallel_rankings"].update({
            "evidence_consensus": artifacts.get("consensus_peptides", ""),
            "event_consensus": artifacts.get("consensus_events", ""),
            "weighted_baseline": artifacts.get("weighted_baseline", ""),
        })

    source_root = result_dir
    consensus_path = Path(artifacts["consensus_peptides"])
    if consensus_path.parent.name == "scoring":
        source_root = consensus_path.parent.parent
    _, validation_rows = read_tsv(artifacts["validation_plan"])
    profile_name = str(
        provenance.get("rules_name")
        or provenance.get("profile")
        or provenance.get("pipeline_version")
        or "evidence_consensus"
    )
    profile = {"_profile_name": profile_name}
    bundle = load_report_bundle(
        profile=profile,
        events=events,
        peptides=report_peptides,
        appm_summary=_summary_mapping(artifacts.get("appm_summary")),
        validation_rows=validation_rows,
        outdir=source_root,
        provenance=provenance,
        sample_id=str(provenance.get("sample_id") or ""),
        entry_mode=str(provenance.get("entry_mode") or ""),
    )
    patient_html = layout.reports / "patient_report.html"
    make_patient_report(patient_html, bundle)
    fusion_html = _fusion_report_html(result_dir)
    if fusion_html:
        current = patient_html.read_text(encoding="utf-8")
        if "融合基因检测结果" not in current:
            current = current.replace("</body></html>", fusion_html + "\n</body></html>")
            patient_html.write_text(current, encoding="utf-8")
    return {"patient_report_html": str(patient_html)}


def _run_appm_ccf_reviews(layout: RunLayout, artifacts: dict[str, str]) -> dict[str, str]:
    appm_dir = layout.review / "hla_loh_appm_review"
    appm_args = ["--outdir", str(appm_dir)]
    for key, flag in (("evidence_report", "--evidence-report"), ("appm_gene_status", "--appm-gene-status"), ("appm_submodule_scores", "--appm-submodule-scores"), ("hla_loh_consensus", "--hla-loh"), ("consensus_peptides", "--ranked-peptides")):
        if artifacts.get(key): appm_args += [flag, artifacts[key]]
    appm_review_main(appm_args)
    ccf_dir = layout.review / "ccf_clonality_review"
    ccf_args = ["--outdir", str(ccf_dir), "--ranked-peptides", artifacts["consensus_peptides"]]
    if artifacts.get("ccf_consensus"): ccf_args += ["--ccf", artifacts["ccf_consensus"]]
    if artifacts.get("purity_tsv"): ccf_args += ["--purity-table", artifacts["purity_tsv"]]
    ccf_review_main(ccf_args)
    return {
        "appm_review": str(appm_dir / "appm_escape_review.md"),
        "appm_driver_defects": str(appm_dir / "appm_driver_defects_summary.tsv"),
        "hla_loh_review": str(appm_dir / "hla_loh_consensus.tsv"),
        "ccf_review": str(ccf_dir / "ccf_clonality_review.md"),
        "ccf_confidence_flags": str(ccf_dir / "ccf_confidence_flags.tsv"),
    }


def _normalize_reports(value: Any) -> set[str]:
    if value is None:
        return {"patient", "technical", "onepage"}
    if isinstance(value, str):
        values = [item.strip().lower() for item in value.split(",") if item.strip()]
    else:
        values = [str(item).strip().lower() for item in value if str(item).strip()]
    selected = set(values)
    if not selected or selected == {"none"}:
        return set()
    allowed = {"patient", "technical", "onepage"}
    unknown = selected - allowed
    if unknown:
        raise ValueError("unsupported reports: " + ",".join(sorted(unknown)))
    return selected


def run_review(args: dict[str, Any]) -> dict[str, Any]:
    result_dir = Path(args.get("result_dir") or args.get("input") or "").resolve()
    case_id = safe_identifier(str(args.get("case_id") or result_dir.name or "REVIEW"))
    layout = RunLayout.create(args.get("outdir") or f"work/open-neo-review/{case_id}")
    top_n = int(args.get("top_n") or 12)
    result = MacroResult("open-neo-review", case_id, new_run_id(case_id, "review"), "review")
    try:
        selected_reports = _normalize_reports(args.get("reports"))
    except ValueError as exc:
        result.blocking_issues.append(FailureCode.INVALID_REPORT_SELECTION.value)
        result.steps.append(MacroStep("00", "report-selection", "BLOCKED", str(exc), failure_code=FailureCode.INVALID_REPORT_SELECTION.value))
        result.finish("BLOCKED").write(layout.skill_result)
        return result.to_dict()
    audit(layout, "open_neo_review.start", "START", result_dir=str(result_dir), top_n=top_n, reports=sorted(selected_reports))
    if layout.root.resolve() == result_dir:
        result.blocking_issues.append(FailureCode.REPORT_BOUNDARY_VIOLATION.value)
        result.steps.append(MacroStep("00", "report-output-boundary", "BLOCKED", "Review outdir must differ from source result_dir", failure_code=FailureCode.REPORT_BOUNDARY_VIOLATION.value))
        result.finish("BLOCKED").write(layout.skill_result); return result.to_dict()

    artifacts = discover_result_artifacts(result_dir)
    integrity = audit_review_inputs(artifacts, layout.review / "integrity")
    integrity_outputs = {"integrity_json": str(layout.review / "integrity/review_integrity.json"), "integrity_checks": str(layout.review / "integrity/review_integrity_checks.tsv"), "integrity_blocking": str(layout.review / "integrity/review_blocking_issues.tsv")}
    result.steps.append(MacroStep("01", "result-integrity-check", integrity["status"], detail="missing=" + ",".join(integrity["missing_artifacts"]), outputs=integrity_outputs))
    result.outputs.update(integrity_outputs)
    if integrity["status"] == "NEEDS_RANKING":
        result.blocking_issues.append(FailureCode.NEEDS_RANKING.value); result.finish("NEEDS_RANKING").write(layout.skill_result); return result.to_dict()
    if integrity["status"] == "BLOCKED":
        result.blocking_issues.append(FailureCode.REVIEW_INTEGRITY_BLOCKED.value); result.finish("BLOCKED").write(layout.skill_result); return result.to_dict()

    _, events = read_tsv(artifacts["consensus_events"]); _, peptides = read_tsv(artifacts["consensus_peptides"]); _, all_tool = read_tsv(artifacts["all_tool_results"])
    review_rows = build_review_rows(events, peptides, all_tool)
    if not review_rows:
        result.blocking_issues.append(FailureCode.EVENT_MAPPING_FAILED.value); result.finish("BLOCKED").write(layout.skill_result); return result.to_dict()
    candidate_review = layout.review / "candidate_review.tsv"; write_tsv(candidate_review, review_rows)
    first_batch = select_first_batch(review_rows, top_n); first_batch_path = layout.review / "first_batch_experiment_set.tsv"; write_tsv(first_batch_path, first_batch)
    completion = [row for row in review_rows if row["review_status"] == "COMPLETE_EVIDENCE"]
    manual = [row for row in review_rows if row["experiment_priority"] == "MANUAL_REVIEW_ONLY"]
    write_tsv(layout.review / "evidence_completion_queue.tsv", completion); write_tsv(layout.review / "manual_review_candidates.tsv", manual)
    result.steps.append(MacroStep("02", "event-level-review", "PASS", detail=f"events={len(review_rows)}; first_batch={len(first_batch)}", outputs={"candidate_review": str(candidate_review), "first_batch": str(first_batch_path)}))

    exp = run_experiment_design({"outdir": str(layout.review / "experiment_design"), "candidate_review": str(candidate_review), "first_batch": str(first_batch_path), "ranked_events": artifacts["consensus_events"], "ranked_peptides": artifacts["consensus_peptides"], "top_n": top_n, "therapy_context": args.get("therapy_context") or "research"})
    result.steps.append(MacroStep("03", "experiment-design", exp.get("status", "PARTIAL"), outputs=exp.get("outputs", {})))
    cmp = run_ranking_compare({"outdir": str(layout.review / "ranking_compare"), "left": artifacts["weighted_baseline"], "left_name": "weighted_baseline", "right": artifacts["consensus_peptides"], "right_name": "evidence_consensus"})
    result.steps.append(MacroStep("04", "weighted-vs-consensus-review", cmp.get("status", "PARTIAL"), outputs=cmp.get("outputs", {})))
    mechanism_outputs = _run_appm_ccf_reviews(layout, artifacts)
    result.steps.append(MacroStep("05", "hla-loh-appm-ccf-review", "PASS", outputs=mechanism_outputs))

    concept_outputs: dict[str, str] = {}
    if selected_reports:
        for concept in ("appm", "ccf", "hla loh", "minigene", "elispot"):
            concept_result = run_concept_explainer({"outdir": str(layout.reports / "concepts" / concept.replace(" ", "_")), "concept": concept, "audience": "patient"})
            concept_outputs[concept.replace(" ", "_")] = concept_result["outputs"]["concept_explanation"]
    context = _read_context(args.get("disease_profile"), args.get("clinical_context"))
    support_outputs = {**mechanism_outputs, "ranking_compare_report": cmp.get("outputs", {}).get("report", "")}
    report_outputs = _write_reports(layout, context, review_rows, first_batch, artifacts, integrity, support_outputs, selected_reports, result_dir)
    if "patient" in selected_reports:
        report_outputs.update(_formal_patient_report(layout, result_dir, context, artifacts, events, peptides, all_tool))
    production_reports: dict[str, str] = {}
    if "technical" in selected_reports:
        technical = run_technical_report({"outdir": str(layout.reports / "production_technical"), "result_dir_or_summary": str(result_dir), "pipeline_manifest": artifacts.get("run_manifest", "")})
        production_reports.update({f"production_technical_{key}": value for key, value in technical.get("outputs", {}).items()})
    report_status = "PASS" if selected_reports else "SKIPPED"
    result.steps.append(MacroStep("06", "reports-and-concept-explanations", report_status, outputs={**report_outputs, **concept_outputs, **production_reports}))

    result.outputs.update({
        "candidate_review": str(candidate_review), "first_batch_experiment_set": str(first_batch_path), "fusion_summary": str(layout.review / "fusion_summary.tsv"),
        "evidence_completion_queue": str(layout.review / "evidence_completion_queue.tsv"), "manual_review_candidates": str(layout.review / "manual_review_candidates.tsv"),
        **{f"experiment_{key}": value for key, value in exp.get("outputs", {}).items()}, **{f"comparison_{key}": value for key, value in cmp.get("outputs", {}).items()},
        **mechanism_outputs, **{f"concept_{key}": value for key, value in concept_outputs.items()}, **report_outputs, **production_reports,
        **{f"pipeline_{key}": value for key, value in artifacts.items()},
    })
    final_status = "PASS_WITH_WARNINGS" if integrity["status"] == "PARTIAL" or completion else "PASS"
    if final_status == "PASS_WITH_WARNINGS": result.warnings.append("Some events or integrity layers require evidence completion; missing evidence was not interpreted as negative")
    result.warnings.append("first_batch_experiment_set is a deterministic research heuristic, not an optimized vaccine or treatment set")
    write_json(layout.run_manifest, {"schema_version": "open-neo-review-manifest-v2", "run_id": result.run_id, "case_id": case_id, "source_result_dir": str(result_dir), "source_run_manifest": artifacts["run_manifest"], "source_artifacts": artifacts, "integrity": integrity, "review_outputs": result.outputs, "top_n": top_n, "reports": sorted(selected_reports), "status": final_status})
    result.outputs["review_manifest"] = str(layout.run_manifest)
    update_case_state(layout, case_id=case_id, current_intent="review", status=final_status, source_result_dir=str(result_dir), outputs=result.outputs)
    audit(layout, "open_neo_review.finish", final_status, candidates=len(review_rows), first_batch=len(first_batch), reports=sorted(selected_reports))
    result.finish(final_status).write(layout.skill_result)
    return result.to_dict()
